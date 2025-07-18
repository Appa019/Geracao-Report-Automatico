import os
import PyPDF2
import docx
from pptx import Presentation
import fitz  # PyMuPDF
from pathlib import Path
import re
from typing import List, Dict, Any

class DocumentProcessor:
    """Classe para processar diferentes tipos de documentos"""
    
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.pptx']
    
    def extract_text(self, file_path: str) -> Dict[str, Any]:
        """
        Extrai texto de diferentes tipos de documentos
        
        Args:
            file_path: Caminho para o arquivo
            
        Returns:
            Dict com texto extraído e metadados
        """
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension == '.pdf':
            return self._extract_pdf_text(file_path)
        elif file_extension == '.docx':
            return self._extract_docx_text(file_path)
        elif file_extension == '.pptx':
            return self._extract_pptx_text(file_path)
        else:
            raise ValueError(f"Formato não suportado: {file_extension}")
    
    def _extract_pdf_text(self, file_path: str) -> Dict[str, Any]:
        """Extrai texto de PDF usando PyMuPDF"""
        try:
            doc = fitz.open(file_path)
            text_content = []
            metadata = {
                'total_pages': len(doc),
                'document_type': 'PDF',
                'has_images': False,
                'has_tables': False
            }
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                
                # Detectar imagens
                if page.get_images():
                    metadata['has_images'] = True
                
                # Detectar possíveis tabelas (busca por padrões de texto estruturado)
                lines = text.split('\n')
                table_indicators = 0
                for line in lines:
                    # Contar linhas com múltiplos espaços ou tabs (indicativo de tabelas)
                    if '\t' in line or '  ' in line.strip():
                        table_indicators += 1
                
                if table_indicators > 3:  # Se há muitas linhas estruturadas
                    metadata['has_tables'] = True
                
                text_content.append({
                    'page': page_num + 1,
                    'text': text,
                    'word_count': len(text.split())
                })
            
            doc.close()
            
            return {
                'content': text_content,
                'metadata': metadata,
                'full_text': ' '.join([page['text'] for page in text_content])
            }
            
        except Exception as e:
            raise Exception(f"Erro ao processar PDF: {str(e)}")
    
    def _extract_docx_text(self, file_path: str) -> Dict[str, Any]:
        """Extrai texto de documento Word"""
        try:
            doc = docx.Document(file_path)
            text_content = []
            
            # Extrair paragrafos
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    text_content.append({
                        'paragraph': i + 1,
                        'text': paragraph.text,
                        'style': paragraph.style.name,
                        'word_count': len(paragraph.text.split())
                    })
            
            # Extrair tabelas
            tables_text = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables_text.append(table_data)
            
            metadata = {
                'total_paragraphs': len(doc.paragraphs),
                'total_tables': len(doc.tables),
                'document_type': 'Word Document',
                'has_images': len(doc.inline_shapes) > 0,
                'has_tables': len(doc.tables) > 0
            }
            
            full_text = ' '.join([p['text'] for p in text_content])
            
            return {
                'content': text_content,
                'tables': tables_text,
                'metadata': metadata,
                'full_text': full_text
            }
            
        except Exception as e:
            raise Exception(f"Erro ao processar Word: {str(e)}")
    
    def _extract_pptx_text(self, file_path: str) -> Dict[str, Any]:
        """Extrai texto de apresentação PowerPoint"""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extrair texto de todas as formas
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                text_content.append({
                    'slide': slide_num + 1,
                    'text': ' '.join(slide_text),
                    'word_count': len(' '.join(slide_text).split())
                })
            
            metadata = {
                'total_slides': len(prs.slides),
                'document_type': 'PowerPoint Presentation',
                'has_images': any(hasattr(shape, 'image') for slide in prs.slides for shape in slide.shapes),
                'has_tables': False  # Simplificado para esta versão
            }
            
            full_text = ' '.join([slide['text'] for slide in text_content])
            
            return {
                'content': text_content,
                'metadata': metadata,
                'full_text': full_text
            }
            
        except Exception as e:
            raise Exception(f"Erro ao processar PowerPoint: {str(e)}")
    
    def create_chunks(self, text_content: Dict[str, Any], max_chunk_size: int = 4000) -> List[str]:
        """
        Divide o texto em chunks adaptados ao tamanho do documento
        
        Args:
            text_content: Conteúdo extraído do documento
            max_chunk_size: Tamanho máximo do chunk em caracteres
            
        Returns:
            Lista de chunks de texto
        """
        full_text = text_content['full_text']
        metadata = text_content['metadata']
        
        # Ajustar tamanho do chunk baseado no tamanho do documento
        if metadata['document_type'] == 'PDF':
            pages = metadata['total_pages']
            if pages > 100:
                max_chunk_size = 6000
            elif pages > 50:
                max_chunk_size = 5000
        
        # Dividir por sentenças primeiro
        sentences = re.split(r'[.!?]+', full_text)
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            # Se adicionar esta sentença não exceder o limite
            if len(current_chunk + sentence) <= max_chunk_size:
                current_chunk += sentence + ". "
            else:
                # Salvar chunk atual e começar novo
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + ". "
        
        # Adicionar último chunk
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def get_document_summary(self, text_content: Dict[str, Any]) -> Dict[str, Any]:
        """
        Cria um resumo básico do documento
        
        Args:
            text_content: Conteúdo extraído do documento
            
        Returns:
            Dict com resumo do documento
        """
        full_text = text_content['full_text']
        metadata = text_content['metadata']
        
        # Estatísticas básicas
        word_count = len(full_text.split())
        char_count = len(full_text)
        
        # Extrair possíveis títulos/cabeçalhos
        lines = full_text.split('\n')
        potential_titles = []
        
        for line in lines[:20]:  # Primeiras 20 linhas
            line = line.strip()
            if line and len(line) < 100 and line.isupper():
                potential_titles.append(line)
        
        return {
            'word_count': word_count,
            'char_count': char_count,
            'estimated_reading_time': word_count // 200,  # Assumindo 200 palavras por minuto
            'potential_titles': potential_titles[:5],
            'document_type': metadata['document_type'],
            'has_images': metadata.get('has_images', False),
            'has_tables': metadata.get('has_tables', False)
        }
